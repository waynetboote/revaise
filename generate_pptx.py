from pptx import Presentation

def create_pptx(summary_text, filename="revaise_presentation.pptx"):
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    title = slide.shapes.title
    title.text = "RevAIse Summary"
    
    textbox = slide.shapes.add_textbox(left=200, top=200, width=400, height=200)
    textbox.text = summary_text

    prs.save(filename)
    return filename