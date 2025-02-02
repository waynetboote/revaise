# ppt_generator.py
from pptx import Presentation
from pptx.util import Inches
from reportlab.lib.pagesizes import letter
from reportlab.pdfgen import canvas
from io import BytesIO

def generate_pdf(summary: str, filename: str = "lesson_plan.pdf") -> BytesIO:
    """Generate a printable PDF lesson plan"""
    buffer = BytesIO()
    c = canvas.Canvas(buffer, pagesize=letter)
    width, height = letter
    
    # Set up styles
    c.setFont("Helvetica-Bold", 16)
    c.drawString(72, height - 72, "Lesson Plan Summary")
    
    c.setFont("Helvetica", 12)
    y_position = height - 100
    for line in summary.split('\n'):
        if y_position < 100:
            c.showPage()
            y_position = height - 72
        c.drawString(72, y_position, line)
        y_position -= 15
    
    c.save()
    buffer.seek(0)
    return buffer

def create_pptx(summary_text, filename="revaise_presentation.pptx"):
    """
    Create a PowerPoint presentation containing the summary text.
    Returns the filename of the saved presentation.
    """
    prs = Presentation()
    
    # Use a slide layout that includes a title placeholder (adjust the index as needed).
    slide_layout = prs.slide_layouts[0]  # typically the title slide
    slide = prs.slides.add_slide(slide_layout)
    
    # Set the slide title.
    if slide.shapes.title:
        slide.shapes.title.text = "RevAIse Summary"
    
    # Add a textbox for the summary.
    left = Inches(1)
    top = Inches(2)
    width = Inches(8)
    height = Inches(2)
    textbox = slide.shapes.add_textbox(left, top, width, height)
    text_frame = textbox.text_frame
    text_frame.text = summary_text
    
    prs.save(filename)
    return filename
